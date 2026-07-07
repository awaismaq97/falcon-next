"""
documents.py — Extract text from uploaded documents for chat context.

The chat send flow is JSON (base64 for images), so documents follow the same
"process on the server, send content in the request" shape: the browser uploads
a file here, we extract plain text, and it POSTs that text back with the next
message (which the send flow injects into the model payload). Raw files are never
stored — only the extracted text, for exactly one turn.

Supported: PDF, Word (.docx), Excel (.xlsx/.xlsm), PowerPoint (.pptx), and any
UTF-8 text format (csv, tsv, txt, md, json, xml, html, yaml, log). Legacy binary
.doc/.xls are rejected with a hint to re-save as the modern format.
"""
from __future__ import annotations

import io
import logging

from fastapi import APIRouter, File, HTTPException, UploadFile

logger = logging.getLogger("falcon.documents")

router = APIRouter(prefix="/documents", tags=["documents"])

MAX_UPLOAD_BYTES = 25 * 1024 * 1024  # 25 MB per file
MAX_TEXT_CHARS = 200_000  # cap extracted text so one file can't blow up context

# Text-like formats we decode directly.
TEXT_EXTS = {
    "txt", "md", "markdown", "csv", "tsv", "json", "xml", "html", "htm",
    "yaml", "yml", "log", "ini", "cfg", "rtf", "tex",
}


def _ext(name: str) -> str:
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")  # try empty password
        except Exception:  # noqa: BLE001
            pass
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            txt = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            txt = ""
        if txt.strip():
            parts.append(txt)
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    import docx  # python-docx

    d = docx.Document(io.BytesIO(data))
    parts = [p.text for p in d.paragraphs if p.text and p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
    try:
        wb.close()
    except Exception:  # noqa: BLE001
        pass
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
    return "\n".join(parts)


def _extract_text_bytes(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract(name: str, ext: str, data: bytes) -> str:
    if ext == "pdf":
        return _extract_pdf(data)
    if ext == "docx":
        return _extract_docx(data)
    if ext in ("xlsx", "xlsm"):
        return _extract_xlsx(data)
    if ext == "pptx":
        return _extract_pptx(data)
    if ext in TEXT_EXTS:
        return _extract_text_bytes(data)
    if ext in ("doc", "xls", "ppt"):
        raise HTTPException(
            status_code=415,
            detail=f"Legacy .{ext} files aren't supported — re-save as .{ext}x and try again.",
        )
    # Last resort: attempt a text decode for unknown extensions.
    text = _extract_text_bytes(data)
    if "\x00" in text[:1000]:  # looks binary
        raise HTTPException(status_code=415, detail=f"Unsupported file type: .{ext or '?'}")
    return text


@router.post("/extract")
async def extract_document(file: UploadFile = File(...)) -> dict:
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file.")
    if len(data) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail="File is larger than 25 MB.")

    name = file.filename or "file"
    ext = _ext(name)
    try:
        text = _extract(name, ext, data)
    except HTTPException:
        raise
    except Exception as exc:  # noqa: BLE001
        logger.warning("extraction failed for %s: %s", name, exc)
        raise HTTPException(status_code=422, detail=f"Could not read {name}: {exc}")

    text = (text or "").strip()
    if not text:
        raise HTTPException(status_code=422, detail=f"No extractable text found in {name}.")

    truncated = len(text) > MAX_TEXT_CHARS
    if truncated:
        text = text[:MAX_TEXT_CHARS]

    return {
        "filename": name,
        "chars": len(text),
        "truncated": truncated,
        "text": text,
    }
